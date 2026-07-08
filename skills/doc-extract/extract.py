#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""doc-extract: unified binary-document -> text extractor for winclaw/metacoder.

Usage:
  python extract.py <file>                                   # single file -> stdout
  python extract.py <file> --out out.txt                     # single file -> file
  python extract.py --batch <dir> --out SOURCE.md            # whole dir -> one merged file
  python extract.py --batch <dir> --out SOURCE.md --max-chars 6000   # cap per file
  python extract.py --batch <dir> --out SOURCE.md --no-ocr   # skip images
Notes (batch defaults, to keep SOURCE.md within the model's context window):
  - de-dupes by stem: if both X.pdf and X.pptx exist, only one is kept (prefer pdf).
  - IMAGES are INCLUDED by default and read via AI vision OCR (see below); pass --no-ocr to skip.
  - --max-chars truncates each file's text (default 6000).

Image AI OCR (PNG/JPG):
  Images can't be Read as text, so they go through an AI vision model (OpenAI-compatible
  /chat/completions with an image_url), extracting all text/tables/labels/figures.
  Config (env, with fallbacks):
    DOC_EXTRACT_VISION_KEY        API key. Fallback: DASHSCOPE_API_KEY, then the qwen key
                                  in winclaw.json (digital-human plugin), then ANTHROPIC_AUTH_TOKEN.
    DOC_EXTRACT_VISION_BASE_URL   OpenAI-compatible base. Default: DashScope compatible-mode
                                  built from WINCLAW_DH_DASHSCOPE_HOST (else dashscope.aliyuncs.com).
    DOC_EXTRACT_VISION_MODEL      Vision model. Default: qwen-vl-max.
  If AI OCR fails, falls back to local tesseract (pytesseract) if installed; otherwise
  emits a clear marker for that one image and continues (never crashes the batch).
"""
import sys, os, argparse, glob, json, base64
import urllib.request


def pdf_to_text(p):
    try:
        import fitz
        return "\n".join(pg.get_text() for pg in fitz.open(p))
    except Exception:
        import pdfplumber
        with pdfplumber.open(p) as pdf:
            return "\n".join((pg.extract_text() or "") for pg in pdf.pages)


def pptx_to_text(p):
    from pptx import Presentation
    out = []
    for i, s in enumerate(Presentation(p).slides, 1):
        out.append("\n--- Slide %d ---" % i)
        for sh in s.shapes:
            if sh.has_text_frame:
                out.append(sh.text_frame.text)
            if getattr(sh, "has_table", False) and sh.has_table:
                for row in sh.table.rows:
                    out.append("\t".join(c.text for c in row.cells))
        if s.has_notes_slide and s.notes_slide.notes_text_frame:
            out.append("[notes] " + s.notes_slide.notes_text_frame.text)
    return "\n".join(out)


def docx_to_text(p):
    import docx
    d = docx.Document(p)
    out = [para.text for para in d.paragraphs]
    for t in d.tables:
        for row in t.rows:
            out.append("\t".join(c.text for c in row.cells))
    return "\n".join(out)


def xlsx_to_text(p):
    import openpyxl
    wb = openpyxl.load_workbook(p, read_only=True, data_only=True)
    out = []
    for ws in wb.worksheets:
        out.append("\n--- Sheet: %s ---" % ws.title)
        for row in ws.iter_rows(values_only=True):
            out.append("\t".join("" if c is None else str(c) for c in row))
    return "\n".join(out)


# ── Image AI OCR ──────────────────────────────────────────────────────────────
_VISION_PROMPT = (
    "提取这张图片中的所有文字、数字、表格、图表(坐标轴/数值/图例)、流程图与标注等"
    "全部信息，用纯文本完整、忠实地输出（保留专有名词、数字、结构与层级），不要解释、"
    "不要总结、不要编造看不到的内容。若图中无文字信息，输出空。"
)


def _vision_key_from_winclaw_json():
    for cand in ("/home/winclaw/.winclaw/winclaw.json",
                 os.path.expanduser("~/.winclaw/winclaw.json")):
        try:
            with open(cand, encoding="utf-8") as f:
                c = json.load(f)
        except Exception:
            continue
        # digital-human plugin's qwen key (DashScope) is the most likely available key
        try:
            dh = (((c.get("plugins") or {}).get("entries") or {}).get("digital-human") or {})
            k = ((dh.get("config") or {}).get("qwen") or {}).get("apiKey")
            if k:
                return k
        except Exception:
            pass
        # any provider apiKey as a last resort
        try:
            for prov in (c.get("models") or {}).get("providers", {}).values():
                if prov.get("apiKey"):
                    return prov["apiKey"]
        except Exception:
            pass
    return None


def _vision_ocr(p):
    key = (os.environ.get("DOC_EXTRACT_VISION_KEY")
           or os.environ.get("DASHSCOPE_API_KEY")
           or _vision_key_from_winclaw_json()
           or os.environ.get("ANTHROPIC_AUTH_TOKEN"))
    if not key:
        raise RuntimeError("no vision API key (set DOC_EXTRACT_VISION_KEY)")
    host = os.environ.get("WINCLAW_DH_DASHSCOPE_HOST", "dashscope.aliyuncs.com")
    base = (os.environ.get("DOC_EXTRACT_VISION_BASE_URL")
            or ("https://%s/compatible-mode/v1" % host))
    model = os.environ.get("DOC_EXTRACT_VISION_MODEL", "qwen-vl-max")
    ext = os.path.splitext(p)[1].lower()
    mime = "image/png" if ext == ".png" else "image/jpeg"
    with open(p, "rb") as f:
        b64 = base64.b64encode(f.read()).decode("ascii")
    payload = {
        "model": model,
        "messages": [{
            "role": "user",
            "content": [
                {"type": "text", "text": _VISION_PROMPT},
                {"type": "image_url",
                 "image_url": {"url": "data:%s;base64,%s" % (mime, b64)}},
            ],
        }],
        "temperature": 0,
    }
    req = urllib.request.Request(
        base.rstrip("/") + "/chat/completions",
        data=json.dumps(payload).encode("utf-8"),
        headers={"Authorization": "Bearer " + key,
                 "Content-Type": "application/json"},
        method="POST",
    )
    with urllib.request.urlopen(req, timeout=180) as r:
        d = json.loads(r.read().decode("utf-8", "replace"))
    return (d["choices"][0]["message"]["content"] or "").strip()


def img_to_text(p):
    # 1) AI vision OCR (primary)
    ai_err = ""
    try:
        t = _vision_ocr(p)
        if t:
            return "[AI-OCR via vision model]\n" + t
        return "[AI-OCR: image had no extractable text]"
    except Exception as e:
        ai_err = str(e)[:200]
    # 2) local tesseract fallback
    try:
        import pytesseract
        from PIL import Image
        t = pytesseract.image_to_string(Image.open(p), lang="jpn+eng+chi_sim")
        return "[tesseract-OCR]\n" + t
    except Exception as e:
        return "[OCR unavailable for %s - AI-OCR(%s) / tesseract(%s)]" % (
            os.path.basename(p), ai_err, str(e)[:120])


DISPATCH = {
    ".pdf": pdf_to_text, ".pptx": pptx_to_text, ".docx": docx_to_text,
    ".xlsx": xlsx_to_text, ".jpg": img_to_text, ".jpeg": img_to_text, ".png": img_to_text,
}
IMG_EXT = (".jpg", ".jpeg", ".png")
TEXT_EXT = (".md", ".txt", ".js", ".json", ".csv", ".html", ".css", ".py", ".ts")
PREFER = {".pdf": 4, ".docx": 3, ".xlsx": 2, ".pptx": 1}  # for stem de-dup


def extract(p):
    ext = os.path.splitext(p)[1].lower()
    fn = DISPATCH.get(ext)
    if fn:
        try:
            return fn(p)
        except Exception as e:
            return "[extract failed %s: %s]" % (p, e)
    try:
        with open(p, encoding="utf-8", errors="replace") as f:
            return f.read()
    except Exception as e:
        return "[read failed %s: %s]" % (p, e)


def select_batch_files(d, ocr):
    cand = [f for f in glob.glob(os.path.join(d, "**", "*"), recursive=True) if os.path.isfile(f)]
    text_files, doc_by_stem = [], {}
    for f in cand:
        base = os.path.basename(f)
        if base.startswith("~$") or base.startswith(".~") or base.lower() == "source.md":
            continue  # skip office temp locks and our own output file
        ext = os.path.splitext(f)[1].lower()
        if ext in TEXT_EXT:
            text_files.append(f)
        elif ext in IMG_EXT:
            if ocr:
                doc_by_stem.setdefault("img:" + f, f)
        elif ext in DISPATCH:  # pdf/pptx/docx/xlsx -> de-dup by stem
            stem = os.path.splitext(os.path.basename(f))[0]
            cur = doc_by_stem.get(stem)
            if cur is None or PREFER.get(ext, 0) > PREFER.get(os.path.splitext(cur)[1].lower(), 0):
                doc_by_stem[stem] = f
    return sorted(set(text_files) | set(doc_by_stem.values()))


def cap(text, n):
    if n and len(text) > n:
        return text[:n] + "\n...[truncated to %d chars]\n" % n
    return text


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("path")
    ap.add_argument("--batch", action="store_true")
    ap.add_argument("--out")
    ap.add_argument("--max-chars", type=int, default=6000, dest="max_chars",
                    help="truncate each file's text (0 = no cap)")
    ap.add_argument("--no-ocr", dest="ocr", action="store_false",
                    help="skip images (default: include images via AI vision OCR)")
    ap.set_defaults(ocr=True)
    a = ap.parse_args()

    if a.batch:
        files = select_batch_files(a.path, a.ocr)
        buf = []
        for f in files:
            buf.append("\n\n# ===== %s =====\n" % os.path.relpath(f, a.path))
            buf.append(cap(extract(f), a.max_chars))
        body = "\n".join(buf)
        text = "# SOURCE (%d files, max %d chars/file, ocr=%s)\n%s" % (
            len(files), a.max_chars, "on" if a.ocr else "off", body)
    else:
        text = cap(extract(a.path), a.max_chars)

    if a.out:
        with open(a.out, "w", encoding="utf-8") as f:
            f.write(text)
        print("[wrote %s - %d chars]" % (a.out, len(text)))
    else:
        sys.stdout.buffer.write(text.encode("utf-8", "replace"))


if __name__ == "__main__":
    main()
